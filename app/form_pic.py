from python_pptx_text_replacer import TextReplacer
from pptx import Presentation
from pptxtopdf import convert



def form_text(res: dict, id):
    replacer = TextReplacer("original.pptx", slides='',
                            tables=True, charts=True, textframes=True)
    replacer.replace_text( [ ('Первый',f'{res['name'].split()[0]}') ] )
    replacer.replace_text( [ ("Второй", f'{res['name'].split()[1]} {res['name'].split()[2]}') ])
    replacer.replace_text( [ ('1902 - 1943', f'{res['era']}') ])
    replacer.replace_text( [ ('Компания', f'{res['rod']}') ])
    path = f"./changed_{id}.pptx"
    replacer.write_presentation_to_file(path)
    return path


def add_photo(img_path, chandged_path):
    root = Presentation(chandged_path)
    
    slide_index = 0
    slide = root.slides[slide_index]

    shape_ind = 4
    shape = slide.shapes[shape_ind]

    picture  = slide.shapes.add_picture(img_path, shape.left, shape.top, shape.width, shape.height)

    root.save(chandged_path) 



def conv_to_pdf(path, id):
    new_path = f'output_{id}'
    convert(path, new_path)
    return f'output_{id}/changed_{id}.pdf'